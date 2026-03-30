"""
EVACYCLE Demo PPT Generator
────────────────────────────────────────────────────────────────
용도: C-level 데모용 PPTX 생성
실행: python docs/make_ppt.py
출력: docs/EVACYCLE-Demo.pptx

의존성: pip install python-pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ─────────────────────────────────────────────
# 색상 팔레트 (EVACYCLE 브랜드 — 시안 A 기반)
# ─────────────────────────────────────────────
C_DARK_BG    = RGBColor(0x0F, 0x17, 0x2A)   # 다크 네이비 (슬라이드 배경)
C_SIDEBAR    = RGBColor(0x1A, 0x2B, 0x3C)   # 사이드 패널
C_GREEN      = RGBColor(0x1A, 0x6B, 0x3C)   # Primary 그린
C_GREEN_LT   = RGBColor(0x22, 0xC5, 0x5E)   # 라이트 그린 (accent)
C_AMBER      = RGBColor(0xF5, 0x9E, 0x0B)   # 앰버 (포인트)
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_GRAY = RGBColor(0xE2, 0xE8, 0xF0)
C_MID_GRAY   = RGBColor(0x94, 0xA3, 0xB8)
C_TEXT_MAIN  = RGBColor(0xF1, 0xF5, 0xF9)
C_TEAL       = RGBColor(0x06, 0xB6, 0xD4)

# 슬라이드 크기 (16:9 와이드)
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ─────────────────────────────────────────────
# 헬퍼 함수
# ─────────────────────────────────────────────

def set_bg(slide, color: RGBColor):
    """슬라이드 배경 단색 채우기."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, fill_color, line_color=None, line_width=None):
    """사각형 도형 추가 → shape 반환."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, l, t, w, h,
                font_size=18, bold=False, color=None,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    """텍스트박스 추가 → text_frame 반환."""
    txBox = slide.shapes.add_textbox(
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color or C_WHITE
    return tf


def add_para(tf, text, font_size=14, bold=False,
             color=None, indent=0, align=PP_ALIGN.LEFT):
    """기존 text_frame에 단락 추가."""
    from pptx.util import Pt
    p = tf.add_paragraph()
    p.alignment = align
    p.level = indent
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color or C_WHITE
    return p


def add_label_value(slide, label, value, l, t, label_w=1.8, total_w=4.5, h=0.35,
                    label_size=10, value_size=13):
    """라벨+값 한 쌍 레이아웃."""
    add_textbox(slide, label, l, t, label_w, h,
                font_size=label_size, color=C_MID_GRAY, bold=False)
    add_textbox(slide, value, l + label_w, t, total_w - label_w, h,
                font_size=value_size, color=C_TEXT_MAIN, bold=True)


def slide_header(slide, title: str, subtitle: str = ""):
    """공통 헤더: 상단 녹색 바 + 제목."""
    # 상단 녹색 바
    add_rect(slide, 0, 0, 13.33, 0.9, C_GREEN)
    # 제목
    add_textbox(slide, title, 0.4, 0.12, 9, 0.7,
                font_size=28, bold=True, color=C_WHITE)
    if subtitle:
        add_textbox(slide, subtitle, 0.4, 0.62, 9, 0.4,
                    font_size=14, color=C_GREEN_LT, bold=False)
    # 우상단 EVACYCLE 워터마크
    add_textbox(slide, "EVACYCLE", 10.8, 0.18, 2.2, 0.5,
                font_size=12, color=RGBColor(0xFF, 0xFF, 0xFF),
                bold=True, align=PP_ALIGN.RIGHT)


def slide_footer(slide, page: int, total: int = 7):
    """슬라이드 번호 하단."""
    add_rect(slide, 0, 7.1, 13.33, 0.4, RGBColor(0x0A, 0x10, 0x1E))
    add_textbox(slide, f"CONFIDENTIAL  ·  C-Level Demo  ·  {page} / {total}",
                0.4, 7.12, 12.5, 0.32,
                font_size=9, color=C_MID_GRAY, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────
# Slide 1 — 표지
# ─────────────────────────────────────────────
def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide, C_DARK_BG)

    # 왼쪽 초록 세로 띠
    add_rect(slide, 0, 0, 0.55, 7.5, C_GREEN)

    # 중앙 녹색 원형 로고 영역 (텍스트로 대체)
    logo_box = add_rect(slide, 2.0, 1.2, 4.0, 4.0, C_GREEN)
    # 로고 글자
    add_textbox(slide, "♻", 3.2, 2.0, 1.8, 1.8,
                font_size=72, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, "EVACYCLE", 2.0, 3.6, 4.0, 0.7,
                font_size=28, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # 제목 영역
    add_textbox(slide, "EVACYCLE", 6.5, 1.5, 6.5, 1.1,
                font_size=52, bold=True, color=C_WHITE)
    add_textbox(slide, "폐차 순환 관리 플랫폼", 6.5, 2.7, 6.5, 0.6,
                font_size=22, bold=False, color=C_GREEN_LT)

    # 구분선
    add_rect(slide, 6.5, 3.4, 6.0, 0.04, C_GREEN)

    # 부가 정보
    tf = add_textbox(slide, "EV 폐차 부품 · 이력관리 · 그레이딩 · 정산 통합 플랫폼", 6.5, 3.6, 6.5, 0.5,
                     font_size=13, color=C_MID_GRAY)
    add_textbox(slide, "v1.0 MVP  ·  2026년 3월", 6.5, 4.2, 6.5, 0.4,
                font_size=11, color=C_AMBER, bold=True)

    slide_footer(slide, 1)


# ─────────────────────────────────────────────
# Slide 2 — 문제 정의
# ─────────────────────────────────────────────
def slide_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_DARK_BG)
    slide_header(slide, "문제 정의", "현행 폐차 프로세스의 구조적 한계")

    problems = [
        ("📋", "추적 불가",
         "폐차 부품의 이동 경로 추적 불가\n종이 기반 체크리스트로 분실·변조 위험"),
        ("💸", "정산 불투명",
         "수기 계산 + 이메일 협의로 정산 지연\n정산 근거 데이터 없음 → 분쟁 빈번"),
        ("📁", "문서 분산",
         "계약서·사진·인수증이 각기 다른 시스템\n감사 추적 불가, 규제 대응 어려움"),
        ("⏱", "처리 지연",
         "역할 간 인계 프로세스 수동화\n허브 입고부터 정산까지 평균 14일"),
    ]

    for i, (icon, title, desc) in enumerate(problems):
        col = i % 2
        row = i // 2
        x = 0.5 + col * 6.3
        y = 1.3 + row * 2.6

        # 카드 배경
        card = add_rect(slide, x, y, 5.9, 2.3, RGBColor(0x1E, 0x2A, 0x3A))
        card.line.color.rgb = C_GREEN
        card.line.width = Pt(1)

        # 아이콘 + 제목
        add_textbox(slide, f"{icon}  {title}", x + 0.2, y + 0.15, 5.5, 0.5,
                    font_size=17, bold=True, color=C_AMBER)
        add_textbox(slide, desc, x + 0.2, y + 0.65, 5.5, 1.5,
                    font_size=12, color=C_LIGHT_GRAY, wrap=True)

    slide_footer(slide, 2)


# ─────────────────────────────────────────────
# Slide 3 — EVACYCLE 솔루션
# ─────────────────────────────────────────────
def slide_solution(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_DARK_BG)
    slide_header(slide, "EVACYCLE 솔루션", "역할별 디지털 플로우 + 블록체인 기반 이력 추적")

    # 플로우 단계
    steps = [
        ("🏭", "폐차장", "케이스 등록\nCoC 서명"),
        ("🏢", "허브", "입고 확인\n그레이딩"),
        ("🏷️", "마켓플레이스", "Lot 등록\n구매 매칭"),
        ("💰", "정산", "M0 + Delta\n자동 계산"),
    ]

    for i, (icon, role, desc) in enumerate(steps):
        x = 0.6 + i * 2.95
        # 원형 아이콘 박스
        box = add_rect(slide, x, 1.2, 2.5, 2.5, C_GREEN)
        box.line.fill.background()
        add_textbox(slide, icon, x, 1.4, 2.5, 0.8,
                    font_size=32, align=PP_ALIGN.CENTER, color=C_WHITE)
        add_textbox(slide, role, x, 2.2, 2.5, 0.5,
                    font_size=15, bold=True, align=PP_ALIGN.CENTER, color=C_WHITE)
        add_textbox(slide, desc, x, 2.7, 2.5, 0.7,
                    font_size=10, align=PP_ALIGN.CENTER, color=C_GREEN_LT)

        # 화살표 (마지막 제외)
        if i < len(steps) - 1:
            add_textbox(slide, "→", x + 2.55, 1.9, 0.4, 0.5,
                        font_size=22, bold=True, color=C_AMBER, align=PP_ALIGN.CENTER)

    # 해시 체인 설명 박스
    add_rect(slide, 0.5, 4.1, 12.3, 1.5, RGBColor(0x0A, 0x2A, 0x1A))
    chain_box = slide.shapes[-1]
    chain_box.line.color.rgb = C_GREEN_LT
    chain_box.line.width = Pt(1)

    add_textbox(slide, "🔗  블록체인 기반 이력 추적 (EventLedger 해시 체인)",
                0.7, 4.2, 8, 0.45, font_size=14, bold=True, color=C_GREEN_LT)
    add_textbox(slide,
                "모든 이벤트(등록·CoC·입고·그레이딩·판매·정산)를 SHA-256 해시 체인으로 연결  ·  위변조 즉시 감지  ·  감사 추적 100%",
                0.7, 4.65, 12.0, 0.7, font_size=11, color=C_LIGHT_GRAY)

    # 핵심 수치 3개
    metrics = [("44개", "API 엔드포인트"), ("6개", "역할 기반 접근제어"), ("100%", "이벤트 감사 추적")]
    for i, (val, lbl) in enumerate(metrics):
        x = 9.5 + i * 1.3
        add_textbox(slide, val, x, 4.2, 1.2, 0.45,
                    font_size=18, bold=True, color=C_AMBER, align=PP_ALIGN.CENTER)
        add_textbox(slide, lbl, x, 4.65, 1.2, 0.4,
                    font_size=9, color=C_MID_GRAY, align=PP_ALIGN.CENTER)

    slide_footer(slide, 3)


# ─────────────────────────────────────────────
# Slide 4 — 시스템 아키텍처
# ─────────────────────────────────────────────
def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_DARK_BG)
    slide_header(slide, "시스템 아키텍처", "역할별 접근제어 + 검증된 기술 스택")

    # 좌측: 역할 구조
    add_textbox(slide, "역할 구조 (RBAC)", 0.4, 1.1, 4.5, 0.4,
                font_size=14, bold=True, color=C_GREEN_LT)

    roles = [
        ("ADMIN",            "플랫폼 운영·정산 승인",   C_AMBER),
        ("JUNKYARD / OWNER", "케이스 등록·CoC 서명",    C_GREEN_LT),
        ("INTAKE_JUNKYARD",  "입고 폐차장 처리",         C_TEAL),
        ("HUB",              "입고 확인·그레이딩·Lot",  C_TEAL),
        ("BUYER",            "마켓플레이스 구매",         C_MID_GRAY),
    ]
    for i, (role, desc, col) in enumerate(roles):
        y = 1.6 + i * 0.72
        add_rect(slide, 0.4, y, 5.5, 0.6, RGBColor(0x1A, 0x28, 0x38))
        # 역할 뱃지
        badge = add_rect(slide, 0.45, y + 0.08, 0.18, 0.44, col)
        add_textbox(slide, role, 0.75, y + 0.1, 2.2, 0.4,
                    font_size=11, bold=True, color=col)
        add_textbox(slide, desc, 2.9, y + 0.12, 2.8, 0.4,
                    font_size=10, color=C_LIGHT_GRAY)

    # 우측: 기술 스택
    add_textbox(slide, "기술 스택", 6.5, 1.1, 6.5, 0.4,
                font_size=14, bold=True, color=C_GREEN_LT)

    stack = [
        ("Frontend",  "Next.js 14 · TypeScript · Tailwind CSS"),
        ("Backend",   "NestJS 10 · TypeScript · REST API"),
        ("Database",  "PostgreSQL 16 · Prisma ORM"),
        ("Cache",     "Redis 7"),
        ("Storage",   "MinIO (S3 호환)"),
        ("Auth",      "OTP(이메일) + JWT (15m/7d)"),
        ("Testing",   "Jest · Supertest · k6"),
        ("CI/CD",     "GitHub Actions → Docker Compose"),
    ]

    for i, (layer, tech) in enumerate(stack):
        y = 1.6 + i * 0.62
        add_rect(slide, 6.5, y, 6.5, 0.55, RGBColor(0x1A, 0x28, 0x38))
        add_textbox(slide, layer, 6.65, y + 0.08, 1.4, 0.4,
                    font_size=10, bold=True, color=C_AMBER)
        add_textbox(slide, tech, 8.1, y + 0.08, 4.7, 0.4,
                    font_size=10, color=C_LIGHT_GRAY)

    slide_footer(slide, 4)


# ─────────────────────────────────────────────
# Slide 5 — 핵심 플로우 (데모 순서)
# ─────────────────────────────────────────────
def slide_demo_flow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_DARK_BG)
    slide_header(slide, "핵심 플로우 (데모 순서)", "DRAFT → SUBMITTED → IN_TRANSIT → RECEIVED → GRADING → ON_SALE → SOLD → SETTLED")

    flow_steps = [
        ("①", "폐차장 로그인", "OTP 인증\n케이스 등록 (차량정보·VIN)", C_GREEN),
        ("②", "CoC 서명", "Chain-of-Custody 디지털 서명\n허브 이동 시작 → IN_TRANSIT", C_TEAL),
        ("③", "허브 입고·그레이딩", "입고 확인 → RECEIVED\n배터리 Grade A 판정 → DerivedLot 생성", C_TEAL),
        ("④", "마켓플레이스 등록", "Lot → Listing 등록 (고정가 ₩1,500,000)\n바이어 카탈로그 노출 → ON_SALE", C_AMBER),
        ("⑤", "바이어 구매", "마켓플레이스 구매 완료 → SOLD\nDelta 정산 자동 생성 (₩225,000)", C_AMBER),
        ("⑥", "관리자 정산 승인", "M0 + Delta 일괄 승인 → PAID\n케이스 자동 SETTLED 전이", C_GREEN_LT),
    ]

    for i, (num, title, desc, col) in enumerate(flow_steps):
        col_idx = i % 3
        row_idx = i // 3
        x = 0.45 + col_idx * 4.25
        y = 1.25 + row_idx * 2.75

        # 카드
        card = add_rect(slide, x, y, 4.0, 2.5, RGBColor(0x14, 0x22, 0x30))
        card.line.color.rgb = col
        card.line.width = Pt(1.5)

        # 번호 원
        num_box = add_rect(slide, x + 0.15, y + 0.15, 0.55, 0.55, col)
        add_textbox(slide, num, x + 0.15, y + 0.14, 0.55, 0.55,
                    font_size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

        # 제목
        add_textbox(slide, title, x + 0.8, y + 0.18, 3.0, 0.45,
                    font_size=14, bold=True, color=col)

        # 구분선
        add_rect(slide, x + 0.15, y + 0.75, 3.7, 0.03, col)

        # 설명
        add_textbox(slide, desc, x + 0.15, y + 0.85, 3.7, 1.5,
                    font_size=11, color=C_LIGHT_GRAY, wrap=True)

    slide_footer(slide, 5)


# ─────────────────────────────────────────────
# Slide 6 — 데모 계정 안내
# ─────────────────────────────────────────────
def slide_accounts(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_DARK_BG)
    slide_header(slide, "데모 계정 안내", "* OTP는 API 서버 콘솔 로그에서 확인")

    # 테이블 헤더
    headers = ["역할", "이메일", "접속 URL", "주요 기능"]
    col_w   = [1.8,    3.5,     3.5,       4.0]
    col_x   = [0.4]
    for w in col_w[:-1]:
        col_x.append(col_x[-1] + w + 0.05)

    header_y = 1.15
    for i, (hdr, w, x) in enumerate(zip(headers, col_w, col_x)):
        add_rect(slide, x, header_y, w, 0.45, C_GREEN)
        add_textbox(slide, hdr, x + 0.1, header_y + 0.06, w - 0.2, 0.35,
                    font_size=12, bold=True, color=C_WHITE)

    # 테이블 행
    rows = [
        ("🛡️ ADMIN",    "admin@evacycle.com",    "localhost:3001/admin",  "대시보드·정산승인·조직관리",  C_AMBER),
        ("🏭 JUNKYARD", "junkyard@evacycle.com", "localhost:3001/cases",  "케이스 등록·CoC 서명",       C_GREEN_LT),
        ("🏢 HUB",      "hub@evacycle.com",      "localhost:3001/lots",   "입고확인·그레이딩·Lot 관리", C_TEAL),
        ("🛒 BUYER",    "buyer@evacycle.com",    "localhost:3001/marketplace", "마켓플레이스·구매",      C_MID_GRAY),
    ]

    row_colors = [RGBColor(0x1A, 0x28, 0x38), RGBColor(0x16, 0x22, 0x30)]
    for r_idx, (role, email, url, feat, col) in enumerate(rows):
        ry = header_y + 0.5 + r_idx * 0.7
        bg = row_colors[r_idx % 2]
        for i, (txt, w, x) in enumerate(zip([role, email, url, feat], col_w, col_x)):
            add_rect(slide, x, ry, w, 0.62, bg)
            fc = col if i == 0 else C_TEXT_MAIN
            add_textbox(slide, txt, x + 0.1, ry + 0.1, w - 0.2, 0.45,
                        font_size=12 if i == 0 else 11, bold=(i == 0), color=fc)

    # 데모 시나리오 박스
    add_rect(slide, 0.4, 4.55, 12.5, 1.8, RGBColor(0x0A, 0x28, 0x18))
    scenario_box = slide.shapes[-1]
    scenario_box.line.color.rgb = C_GREEN
    scenario_box.line.width = Pt(1)

    add_textbox(slide, "📋  데모 시나리오 (권장 순서)", 0.6, 4.62, 5, 0.4,
                font_size=13, bold=True, color=C_GREEN_LT)
    scenarios = [
        "① 폐차장 로그인 → 케이스 목록 확인 (DRAFT/SUBMITTED/IN_TRANSIT 3건 사전 세팅됨)",
        "② 바이어 로그인 → 마켓플레이스 → LOT 구매 (₩1,500,000)",
        "③ 관리자 로그인 → 대시보드 → 정산 일괄 승인 → SETTLED 전이 확인",
    ]
    tf = add_textbox(slide, "", 0.6, 5.1, 12.0, 1.15,
                     font_size=11, color=C_LIGHT_GRAY)
    tf.paragraphs[0].runs[0].text = scenarios[0]
    for s in scenarios[1:]:
        add_para(tf, s, font_size=11, color=C_LIGHT_GRAY)

    slide_footer(slide, 6)


# ─────────────────────────────────────────────
# Slide 7 — 로드맵
# ─────────────────────────────────────────────
def slide_roadmap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_DARK_BG)
    slide_header(slide, "로드맵", "백엔드 MVP 완료 → 프론트엔드 → AWS 배포")

    # 타임라인 라인
    add_rect(slide, 0.8, 3.7, 11.8, 0.06, C_GREEN)

    milestones = [
        ("v1.0",  "2026.03",  "백엔드 MVP\n완료 ✅",
         "NestJS API 44개\nFlow A E2E 검증\nVera QA 통과",
         C_GREEN_LT, True),
        ("v1.4",  "2026.04.01", "프론트엔드\n완료 예정",
         "Next.js 웹앱\n시안 A/B 역할별\n분리 적용",
         C_AMBER, False),
        ("v2.0",  "2026.Q2",  "AWS 배포\n예정",
         "ECS Fargate\nRDS + ElastiCache\nCI/CD 파이프라인",
         C_TEAL, False),
    ]

    for i, (ver, date, title, detail, col, done) in enumerate(milestones):
        x = 1.5 + i * 3.8

        # 타임라인 원 마커
        marker = add_rect(slide, x + 0.85, 3.38, 0.7, 0.7,
                          col if done else RGBColor(0x1A, 0x28, 0x38))
        if not done:
            marker.line.color.rgb = col
            marker.line.width = Pt(2)

        add_textbox(slide, "●" if done else "○",
                    x + 0.85, 3.38, 0.7, 0.7,
                    font_size=28, color=col, align=PP_ALIGN.CENTER)

        # 버전 + 날짜 (위)
        add_textbox(slide, ver, x + 0.6, 2.55, 1.2, 0.4,
                    font_size=16, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_textbox(slide, date, x + 0.3, 2.95, 1.8, 0.38,
                    font_size=10, color=C_MID_GRAY, align=PP_ALIGN.CENTER)

        # 카드 (아래)
        card = add_rect(slide, x + 0.1, 4.25, 3.3, 2.5,
                        RGBColor(0x14, 0x22, 0x30))
        card.line.color.rgb = col
        card.line.width = Pt(1 if not done else 1.5)

        # DONE 배지
        if done:
            badge = add_rect(slide, x + 0.2, 4.3, 0.9, 0.35, col)
            add_textbox(slide, "DONE ✅", x + 0.2, 4.3, 0.9, 0.35,
                        font_size=8, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

        add_textbox(slide, title, x + 0.2, 4.72, 3.0, 0.55,
                    font_size=13, bold=True, color=col)
        add_textbox(slide, detail, x + 0.2, 5.3, 3.0, 1.3,
                    font_size=10, color=C_LIGHT_GRAY, wrap=True)

    # 하단 문구
    add_textbox(slide,
                "💡  프론트엔드 설계 완료 (시안 A·B 역할 분리)  ·  Finn 구현 착수 예정",
                1.5, 6.65, 10.3, 0.4,
                font_size=11, color=C_AMBER, align=PP_ALIGN.CENTER, bold=True)

    slide_footer(slide, 7)


# ─────────────────────────────────────────────
# MAIN
# ─────────────────────────────────────────────
def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    print("🖥️  EVACYCLE Demo PPT 생성 중...")

    slide_cover(prs)        ; print("  ✅ Slide 1 — 표지")
    slide_problem(prs)      ; print("  ✅ Slide 2 — 문제 정의")
    slide_solution(prs)     ; print("  ✅ Slide 3 — 솔루션")
    slide_architecture(prs) ; print("  ✅ Slide 4 — 아키텍처")
    slide_demo_flow(prs)    ; print("  ✅ Slide 5 — 데모 플로우")
    slide_accounts(prs)     ; print("  ✅ Slide 6 — 계정 안내")
    slide_roadmap(prs)      ; print("  ✅ Slide 7 — 로드맵")

    out_dir = os.path.dirname(os.path.abspath(__file__))
    out_path = os.path.join(out_dir, "EVACYCLE-Demo.pptx")
    prs.save(out_path)
    print(f"\n🎉  저장 완료: {out_path}")
    print("    슬라이드: 7장  |  포맷: 16:9 와이드")


if __name__ == "__main__":
    main()
